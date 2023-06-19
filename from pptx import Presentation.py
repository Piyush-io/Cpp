from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Function to add a new slide with a title and content
def add_slide_with_content(title, content):
    slide_layout = presentation.slide_layouts[1]  # Choose a slide layout (1 for title and content)
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    content_placeholder = slide.placeholders[1]
    content_placeholder.text = content

# Function to add a new slide with an image
def add_slide_with_image(image_path):
    slide_layout = presentation.slide_layouts[5]  # Choose a slide layout (5 for picture with caption)
    slide = presentation.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1)
    slide.shapes.add_picture(image_path, left, top)

# Function to add a new slide with a colored background
def add_slide_with_color_background(color):
    slide_layout = presentation.slide_layouts[5]  # Choose a slide layout (5 for picture with caption)
    slide = presentation.slides.add_slide(slide_layout)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])

# Add slides with content, images, and colored backgrounds
add_slide_with_content("Title Slide", "This is the content of the first slide.")
add_slide_with_image("path_to_image.jpg")
add_slide_with_color_background((255, 0, 0))  # Example: Red color background

# Save the presentation
presentation.save("output.pptx")
